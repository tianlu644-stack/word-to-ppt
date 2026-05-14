#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Word 转 PowerPoint 转换工具 - 高端4A风格

功能：
- 读取 Word 文档内容
- 按照高端4A广告公司调性生成 PowerPoint 幻灯片
- 深色粉红主题配色
- 专业排版和装饰
"""

import os
import re
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


class AdvancedPPTConverter:
    """高端4A风格的 Word 到 PPT 转换器"""
    
    # 深色粉红主色系
    COLOR_DARK_PINK = RGBColor(180, 50, 120)      # 深粉红主色
    COLOR_LIGHT_PINK = RGBColor(230, 150, 180)    # 浅粉红辅助色
    COLOR_ACCENT = RGBColor(255, 100, 150)        # 亮粉红强调色
    COLOR_DARK_BG = RGBColor(25, 25, 35)          # 深灰黑背景
    COLOR_WHITE = RGBColor(255, 255, 255)         # 白色文字
    COLOR_LIGHT_GRAY = RGBColor(200, 200, 210)    # 浅灰文字
    
    def __init__(self, word_file):
        """
        初始化转换器
        
        Args:
            word_file (str): Word 文件路径
        """
        self.word_file = word_file
        self.doc = Document(word_file)
        self.prs = Presentation()
        
        # 设置 16:9 宽屏
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        
    def _add_dark_background(self, slide):
        """添加深色背景"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK_BG
    
    def _add_decorative_line(self, slide, y_position, color=None, width=Pt(3)):
        """添加装饰线条"""
        if color is None:
            color = self.COLOR_DARK_PINK
        
        left = Inches(0.8)
        width_pixels = Inches(3.5)
        line = slide.shapes.add_connector(1, left, y_position, left + width_pixels, y_position)
        line.line.color.rgb = color
        line.line.width = width
    
    def _add_accent_shape(self, slide, left, top, width, height):
        """添加装饰几何形状"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.COLOR_ACCENT
        shape.line.color.rgb = self.COLOR_DARK_PINK
        shape.line.width = Pt(2)
        return shape
    
    def add_title_slide(self, title, subtitle=""):
        """
        添加标题幻灯片 - 高端大气
        
        Args:
            title (str): 标题
            subtitle (str): 副标题
        """
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 深色背景
        self._add_dark_background(slide)
        
        # 添加大背景装饰形状
        accent_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8), Inches(-0.5), Inches(6), Inches(9)
        )
        accent_bg.fill.solid()
        accent_bg.fill.fore_color.rgb = self.COLOR_DARK_PINK
        accent_bg.fill.transparency = 0.85  # 半透明
        accent_bg.line.color.rgb = self.COLOR_DARK_PINK
        accent_bg.line.width = Pt(0)
        
        # 标题文字框
        left = Inches(1)
        top = Inches(2)
        width = Inches(7)
        height = Inches(2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 标题格式
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(72)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.COLOR_WHITE
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.line_spacing = 1.2
        
        # 添加强调线
        self._add_decorative_line(slide, Inches(1.9), self.COLOR_ACCENT, Pt(4))
        
        # 副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(left, top + Inches(2.2), width, Inches(1.5))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.word_wrap = True
            
            for paragraph in subtitle_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.color.rgb = self.COLOR_LIGHT_PINK
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.space_before = Pt(12)
    
    def add_content_slide(self, title, content_list):
        """
        添加内容幻灯片 - 专业版面设计
        
        Args:
            title (str): 标题
            content_list (list): 内容列表
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 深色背景
        self._add_dark_background(slide)
        
        # 左侧装饰竖条
        accent_left = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        accent_left.fill.solid()
        accent_left.fill.fore_color.rgb = self.COLOR_ACCENT
        accent_left.line.width = Pt(0)
        
        # 标题区域背景
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.4), Inches(12.333), Inches(1)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = self.COLOR_DARK_PINK
        title_bg.fill.transparency = 0.3
        title_bg.line.width = Pt(0)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(48)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.COLOR_ACCENT
        
        # 添加强调线
        self._add_decorative_line(slide, Inches(1.5), self.COLOR_ACCENT, Pt(3))
        
        # 内容区域
        content_top = Inches(1.8)
        content_box = slide.shapes.add_textbox(
            Inches(0.8), content_top, Inches(11.733), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i > 0:
                text_frame.add_paragraph()
            
            p = text_frame.paragraphs[i]
            p.text = f"• {item}"
            p.font.size = Pt(26)
            p.font.color.rgb = self.COLOR_WHITE
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            p.level = 0
            
            # 首行缩进
            p.paragraph_format.left_indent = Inches(0.3)
            p.paragraph_format.first_line_indent = Inches(-0.3)
    
    def add_section_slide(self, section_title):
        """
        添加分章节幻灯片 - 过渡页面
        
        Args:
            section_title (str): 章节标题
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 深色背景
        self._add_dark_background(slide)
        
        # 大背景装饰
        accent_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6), Inches(1), Inches(7.333), Inches(6.5)
        )
        accent_bg.fill.solid()
        accent_bg.fill.fore_color.rgb = self.COLOR_DARK_PINK
        accent_bg.fill.transparency = 0.8
        accent_bg.line.width = Pt(0)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(3))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(66)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.COLOR_ACCENT
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.line_spacing = 1.1
        
        # 底部装饰线
        bottom_line = slide.shapes.add_connector(
            1, Inches(1.5), Inches(6.5), Inches(8), Inches(6.5)
        )
        bottom_line.line.color.rgb = self.COLOR_ACCENT
        bottom_line.line.width = Pt(3)
    
    def add_closing_slide(self, title="感谢观看"):
        """
        添加结束幻灯片
        
        Args:
            title (str): 结束文案
        """
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 深色背景
        self._add_dark_background(slide)
        
        # 大型装饰元素
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(7), Inches(1), Inches(5), Inches(5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = self.COLOR_DARK_PINK
        circle.fill.transparency = 0.7
        circle.line.width = Pt(0)
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(72)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER
    
    def _is_title_like(self, text):
        """判断文本是否像标题"""
        # 短文本或全是数字/符号的被认为是标题
        return len(text) < 40 and len(text) > 0
    
    def _split_into_sections(self, paragraphs):
        """将段落分组成多个章节"""
        sections = []
        current_section = {
            'title': '',
            'content': []
        }
        
        for para in paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # 如果是短文本且已有内容，则认为是新标题
            if self._is_title_like(text) and current_section['content']:
                if current_section['title']:
                    sections.append(current_section)
                current_section = {
                    'title': text,
                    'content': []
                }
            else:
                if not current_section['title']:
                    current_section['title'] = text if self._is_title_like(text) else '主要���容'
                else:
                    current_section['content'].append(text)
        
        # 添加最后一个章节
        if current_section['title'] or current_section['content']:
            sections.append(current_section)
        
        return sections
    
    def convert(self):
        """
        从 Word 文档转换为 PPT
        
        Returns:
            Presentation: 生成的演示文稿
        """
        paragraphs = self.doc.paragraphs
        
        if not paragraphs:
            print("Word 文档为空")
            return self.prs
        
        # 第一个段落作为主标题
        first_text = paragraphs[0].text.strip()
        if first_text:
            self.add_title_slide(first_text)
        
        # 分组处理内容
        sections = self._split_into_sections(paragraphs[1:])
        
        for i, section in enumerate(sections):
            title = section['title']
            content = section['content']
            
            # 如果有内容，先添加分章节页面
            if content and i > 0:
                self.add_section_slide(title)
            
            # 添加内容页面
            if content:
                self.add_content_slide(title, content)
            elif title != first_text:
                self.add_section_slide(title)
        
        # 添加结束页面
        self.add_closing_slide("感谢观看")
        
        return self.prs
    
    def save(self, output_file):
        """
        保存 PPT 文件
        
        Args:
            output_file (str): 输出文件路径
        """
        self.prs.save(output_file)
        print(f"✓ PPT 已保存到: {output_file}")


def main():
    """主函数"""
    import sys
    
    if len(sys.argv) < 2:
        print("使用方法: python word_to_ppt.py <Word文件路径> [输出PPT路径]")
        print("示例: python word_to_ppt.py document.docx output.pptx")
        return
    
    word_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else word_file.replace('.docx', '.pptx')
    
    if not os.path.exists(word_file):
        print(f"错误: 文件 '{word_file}' 不存在")
        return
    
    print(f"正在转换: {word_file}")
    print("风格: 高端4A广告公司 - 深色粉红主题")
    
    try:
        converter = AdvancedPPTConverter(word_file)
        converter.convert()
        converter.save(output_file)
        print(f"✓ 转换成功！")
    except Exception as e:
        print(f"✗ 转换失败: {e}")
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    main()
